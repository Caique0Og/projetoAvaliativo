import os
import sys
from pptx import Presentation
import subprocess

def extrair_perguntas(pptx_file):
    """Extrai as perguntas e seus números de um arquivo PPTX."""
    prs = Presentation(pptx_file)
    perguntas = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = row.cells
                    # Verifica se há pelo menos duas células na linha
                    if len(cells) >= 2:
                        try:
                            # Extrai o texto da primeira célula
                            numero_cell_text = cells[0].text.strip()
                            # Remove espaços em branco e caracteres não numéricos
                            numero_str = ''.join(c for c in numero_cell_text if c.isdigit())

                            # Extrai o texto da segunda célula
                            texto_pergunta = cells[1].text.strip()

                            # Verifica se o número da pergunta foi encontrado
                            if numero_str:
                                numero = int(numero_str)
                                perguntas.append((numero, texto_pergunta))
                            else:
                                print("Número da pergunta não encontrado ou vazio.")

                        except ValueError:
                            # Ignora se a primeira célula não contiver um número inteiro
                            pass
                        except Exception as e:
                            print(f"Erro ao processar linha: {e}")

    return perguntas

def apresentar_questionario(perguntas):
    """Apresenta o questionário ao usuário e coleta as respostas."""
    respostas = {}
    for numero, pergunta in perguntas:
        while True:
            resposta = input(f"Pergunta {numero}: {pergunta} (0/1): ").strip()
            if resposta in ['0', '1']:
                respostas[numero] = int(resposta)
                break
            else:
                print("Resposta inválida. Por favor, responda '0' ou '1'.")
    return respostas

def calcular_pontuacao(respostas):
    """Calcula a pontuação com base nas respostas."""
    pontuacao = sum(respostas.values())
    return pontuacao

def fornecer_feedback(pontuacao, total_perguntas):
    """Fornece feedback com base na pontuação."""
    percentual = (pontuacao / total_perguntas) * 100
    print(f"\nPontuação total: {pontuacao} de {total_perguntas} ({percentual:.2f}%)")

    if percentual >= 75:
        print("Sua startup está em um estágio avançado de maturidade. Parabéns!")
    elif percentual >= 50:
        print("Sua startup está em um estágio intermediário de maturidade. Continue trabalhando para melhorar.")
    else:
        print("Sua startup está em um estágio inicial de maturidade. Concentre-se em fortalecer as áreas mais fracas.")

def salvar_resultados(respostas, filename="resultados.txt"):
    """Salva os resultados em um arquivo."""
    with open(filename, 'w') as f:
        for numero, resposta in respostas.items():
            f.write(f"Pergunta {numero}: {resposta}\n")
    print(f"Resultados salvos em {filename}")

def abrir_arquivo(filename):
    """Abre o arquivo usando o programa padrão do sistema."""
    try:
        if sys.platform.startswith('darwin'):  # macOS
            subprocess.call(('open', filename))
        elif os.name == 'nt':  # Windows
            os.startfile(filename)
        elif os.name == 'posix':  # Linux
            subprocess.call(('xdg-open', filename))
        else:
            print("Sistema operacional não suportado para abrir o arquivo automaticamente.")
    except FileNotFoundError:
        print(f"Erro: O arquivo '{filename}' não foi encontrado.")
    except Exception as e:
        print(f"Erro ao abrir o arquivo: {e}")

def colorir_numero(numero, resposta):
    """Retorna o número formatado com a cor apropriada."""
    if resposta == 0:
        return f"\033[91m{numero}\033[0m"  # Vermelho
    else:
        return f"\033[92m{numero}\033[0m"  # Verde

def main():
    """Função principal para executar o questionário."""
    # Obtenha o caminho absoluto para o diretório do script
    script_dir = os.path.dirname(os.path.abspath(__file__))
    # Construa o caminho absoluto para o arquivo PPTX
    pptx_file = os.path.join(script_dir, "Cenario-2-000000000000000000000000000001.pptx")  # Substitua pelo caminho correto

    # Verifique se o arquivo existe
    if not os.path.exists(pptx_file):
        print(f"Erro: O arquivo '{pptx_file}' não foi encontrado.")
        return
    
    # Abre o arquivo PPTX
    abrir_arquivo(pptx_file)
    
    perguntas = extrair_perguntas(pptx_file)

    if not perguntas:
        print("Nenhuma pergunta encontrada no arquivo PPTX.")
        return

    respostas = apresentar_questionario(perguntas)

    print("\nRespostas:")
    for numero, resposta in respostas.items():
        numero_colorido = colorir_numero(numero, resposta)
        print(f"Pergunta {numero_colorido}: {resposta}")

    pontuacao = calcular_pontuacao(respostas)
    fornecer_feedback(pontuacao, len(perguntas))
    salvar_resultados(respostas)

if __name__ == "__main__":
    main()
